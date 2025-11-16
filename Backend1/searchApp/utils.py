# utils.py
import os
import shutil
import logging

from whoosh import index
from whoosh.fields import Schema, TEXT, ID, KEYWORD
from whoosh.qparser import MultifieldParser
import pdfplumber
import docx
from pptx import Presentation
import openpyxl

logger = logging.getLogger(__name__)

# Directory for whoosh index (ensure exists)
WHOOSH_DIR = os.path.join(os.path.dirname(__file__), "whoosh_index")
os.makedirs(WHOOSH_DIR, exist_ok=True)

# Keyword rules for auto-tagging
KEYWORD_RULES = {
    "Marketing": ["marketing", "campaign", "social media", "ads", "email campaign"],
    "Design": ["design", "creative", "illustration", "figma", "photoshop", "ux", "ui"],
    "Product": ["release", "product", "launch", "feature", "roadmap"],
    "SEO": ["seo", "keyword", "backlink", "organic"],
    "Sales": ["pricing", "discount", "sale", "offer"]
}

def _schema():
    return Schema(
        id=ID(stored=True, unique=True),
        title=TEXT(stored=True),
        content=TEXT(stored=True),
        tags=KEYWORD(stored=True, commas=True)
    )

def ensure_index():
    """
    Create/open whoosh index. Raises if creation/open fails.
    """
    try:
        os.makedirs(WHOOSH_DIR, exist_ok=True)
    except Exception as e:
        logger.exception("Failed to create WHOOSH_DIR %s: %s", WHOOSH_DIR, e)
        raise

    try:
        if index.exists_in(WHOOSH_DIR):
            ix = index.open_dir(WHOOSH_DIR)
            logger.debug("Whoosh index opened at %s", WHOOSH_DIR)
            return ix
        else:
            ix = index.create_in(WHOOSH_DIR, _schema())
            logger.info("Whoosh index created at %s", WHOOSH_DIR)
            return ix
    except Exception as e:
        logger.exception("Failed to open/create whoosh index at %s: %s", WHOOSH_DIR, e)
        raise

def extract_text_from_file(path):
    """
    Attempt to extract text for common types. Returns string.
    """
    path = str(path)
    low = path.lower()

    if not os.path.exists(path):
        logger.debug("extract_text_from_file: file not found: %s", path)
        return ""

    # PDF
    if low.endswith(".pdf"):
        try:
            text = []
            with pdfplumber.open(path) as pdf:
                for p in pdf.pages:
                    text.append(p.extract_text() or "")
            return "\n".join(text)
        except Exception as e:
            logger.warning("PDF parse failed %s: %s", path, e)
            return ""

    # DOCX
    if low.endswith(".docx"):
        try:
            doc = docx.Document(path)
            return "\n".join([p.text for p in doc.paragraphs])
        except Exception as e:
            logger.warning("DOCX parse failed %s: %s", path, e)
            return ""

    # PPTX
    if low.endswith(".pptx"):
        try:
            prs = Presentation(path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                    except Exception:
                        continue
            return "\n".join(text_parts)
        except Exception as e:
            logger.warning("PPTX parse failed %s: %s", path, e)
            return ""

    # XLSX / XLSM
    if low.endswith(".xlsx") or low.endswith(".xlsm"):
        try:
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            text_parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            text_parts.append(str(cell))
            return "\n".join(text_parts)
        except Exception as e:
            logger.warning("XLSX parse failed %s: %s", path, e)
            return ""

    # TXT fallback
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        logger.debug("Text read failed %s: %s", path, e)
        return ""

def auto_tag(text):
    if not text:
        return ""
    text_low = text.lower()
    tags = set()
    for tag, keywords in KEYWORD_RULES.items():
        for kw in keywords:
            if kw in text_low:
                tags.add(tag)
                break
    return ",".join(sorted(tags))

def index_document(doc):
    """
    Index/update a single Document model instance.
    """
    try:
        ix = ensure_index()
    except Exception as e:
        logger.exception("index_document: cannot access/create index: %s", e)
        return

    try:
        writer = ix.writer()
        writer.update_document(
            id=str(doc.id),
            title=doc.title or "",
            content=doc.content or "",
            tags=(doc.tags or "")
        )
        writer.commit()
        logger.debug("Indexed document id=%s title=%s", getattr(doc, "id", "<no id>"), getattr(doc, "title", ""))
    except Exception as e:
        logger.exception("Error indexing document %s: %s", getattr(doc, "id", "<no id>"), e)

def delete_document_index(doc_id):
    """
    Remove a document from the whoosh index by id (string or int).
    """
    try:
        ix = ensure_index()
    except Exception as e:
        logger.exception("delete_document_index: cannot access index: %s", e)
        return

    try:
        writer = ix.writer()
        writer.delete_by_term("id", str(doc_id))
        writer.commit()
        logger.debug("Deleted index entry for id=%s", doc_id)
    except Exception as e:
        logger.exception("Error deleting document from index %s: %s", doc_id, e)

def search_query(q, limit=30):
    """
    Search whoosh index for query `q`. Returns list of dicts:
    [{id, title, snippet, tags}, ...]
    Returns [] on any failure (caller can fallback).
    """
    try:
        ix = ensure_index()
    except Exception as e:
        logger.exception("Cannot access whoosh index; search skipped: %s", e)
        return []

    try:
        parser = MultifieldParser(["title", "content", "tags"], ix.schema)
        try:
            qobj = parser.parse(q)
        except Exception as parse_err:
            logger.warning("Whoosh parse failed for query '%s': %s. Trying quoted fallback.", q, parse_err)
            try:
                qobj = parser.parse('"%s"' % q.replace('"', ''))
            except Exception as e2:
                logger.exception("Fallback parse also failed for '%s': %s", q, e2)
                return []

        out = []
        with ix.searcher() as s:
            results = s.search(qobj, limit=limit)
            # allow long highlights
            try:
                results.fragmenter.charlimit = None
            except Exception:
                pass

            for r in results:
                snippet = r.highlights("content") or ""
                out.append({
                    "id": r["id"],
                    "title": r["title"],
                    "snippet": snippet,
                    "tags": r.get("tags", "")
                })

        logger.debug("Whoosh search for '%s' returned %d results", q, len(out))
        return out
    except Exception as e:
        logger.exception("Search failed for query %s: %s", q, e)
        return []

def reindex_all_documents():
    """
    Rebuild the whoosh index for all documents.
    Usage (from manage.py shell):
        from searchapp import utils
        utils.reindex_all_documents()
    """
    from .models import Document

    # remove existing index dir (safest simple approach)
    try:
        if os.path.exists(WHOOSH_DIR):
            shutil.rmtree(WHOOSH_DIR)
            logger.info("Removed existing whoosh index dir: %s", WHOOSH_DIR)
    except Exception as e:
        logger.exception("Failed to remove whoosh dir %s: %s", WHOOSH_DIR, e)
        return

    try:
        ix = index.create_in(WHOOSH_DIR, _schema())
    except Exception as e:
        logger.exception("Failed to create whoosh index during reindex: %s", e)
        return

    writer = ix.writer()
    count = 0
    for doc in Document.objects.all():
        try:
            writer.update_document(id=str(doc.id), title=doc.title or "", content=doc.content or "", tags=(doc.tags or ""))
            count += 1
        except Exception:
            logger.exception("Failed adding doc %s to index", getattr(doc, "id", "<no id>"))
    try:
        writer.commit()
        logger.info("Reindexed %d documents into whoosh index", count)
    except Exception as e:
        logger.exception("Failed to commit reindex: %s", e)
